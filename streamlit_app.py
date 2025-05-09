import streamlit as st
import pandas as pd
import plotly.graph_objs as go
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import os
from datetime import datetime
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HCAHPS_PATH = 'HCAHPS.csv'
HOSPITAL_INFO_PATH = 'Hospital_General_Information.csv'

@st.cache_data
def load_hospital_info():
    try:
        st.write("Loading hospital information...")
        df = pd.read_csv("https://www.dropbox.com/scl/fi/fq5o8a6evwpsfzutjp7uw/Hospital_General_Information.csv?rlkey=c60s0se15d6nzs40mm19a2q5v&st=li48t6ft&dl=1", dtype=str)
        st.write("Hospital information loaded successfully!")
        return df
    except Exception as e:
        st.error(f"Error loading hospital information: {str(e)}")
        return pd.DataFrame()

@st.cache_data
def load_hcahps_data():
    try:
        st.write("Loading HCAHPS data...")
        df = pd.read_csv("https://www.dropbox.com/scl/fi/d35e3po3qfyaw7fz3qend/HCAHPS.csv?rlkey=pw76uj8z5270ks7izz6esx62r&st=ugsr5p6s&dl=1", dtype=str)
        df.columns = df.columns.str.strip()
        df['HCAHPS Answer Percent'] = pd.to_numeric(df['HCAHPS Answer Percent'], errors='coerce')
        st.write("HCAHPS data loaded successfully!")
        return df
    except Exception as e:
        st.error(f"Error loading HCAHPS data: {str(e)}")
        return pd.DataFrame()

st.title('🏥 Hospital HCAHPS Benchmarking Tool')

with st.expander('ℹ️ What do these metrics mean? (Click to expand)', expanded=False):
    st.markdown('''
**How to interpret the table and chart:**
- **Hospital**: The selected hospital's average score for each metric (higher is better).
- **State Avg**: The average score for all hospitals in the same state.
- **National Avg**: The average score for all hospitals in the U.S.
- **vs State / vs National**: How much higher or lower the hospital is compared to the state/national average (green = better, red = worse).

**Metric definitions:**
- **Nurse Communication**: % of patients who said nurses "always" communicated well.
- **Doctor Communication**: % of patients who said doctors "always" communicated well.
- **Staff Responsiveness**: % of patients who said they "always" received help as soon as they wanted.
- **Care Transition**: % of patients who "strongly agree" they understood their care when leaving the hospital.
- **Discharge Info**: % of patients who said staff "did" give them discharge information.
- **Care Cleanliness**: % of patients who said their room was "always" clean.
- **Quietness**: % of patients who said the area around their room was "always" quiet at night.
- **Recommend**: % of patients who would "definitely recommend" the hospital.

**How to use this:**
- Look for green cells: These are areas where the hospital outperforms the benchmark.
- Look for red cells: These are areas where the hospital underperforms and may need improvement.
- Use the chart to quickly spot strengths and weaknesses compared to state and national averages.
''')

df_info = load_hospital_info()
df_hcahps = load_hcahps_data()
# st.write("All unique HCAHPS Answer Descriptions:", df_hcahps['HCAHPS Answer Description'].unique())
hospital_names = df_info['Facility Name'].sort_values().unique()
hospital = st.selectbox('Select Hospital', hospital_names)

hospital_row = df_info[df_info['Facility Name'] == hospital].iloc[0]
hospital_id = hospital_row['Facility ID']
hospital_state = hospital_row['State']

# Define the measures and the answer description to use for each
all_measures = [
    ("Nurse Communication", "H_COMP_1_A_P"),
    ("Doctor Communication", "H_COMP_2_A_P"),
    ("Staff Responsiveness", "H_COMP_3_A_P"),
    ("Care Transition", "H_COMP_5_A_P"),
    ("Discharge Info", "H_COMP_6_Y_P"),
    ("Care Cleanliness", "H_CLEAN_HSP_A_P"),
    ("Quietness", "H_QUIET_HSP_A_P"),
    ("Recommend", "H_RECMND_DY"),
]

metric_options = [m[0] for m in all_measures]
selected_metrics = st.multiselect(
    "Select metrics to display", metric_options, default=metric_options
)
measures = [m for m in all_measures if m[0] in selected_metrics]

def match_answer(series, answer):
    return series.str.strip().str.lower() == answer.strip().lower()

comparison = []
for label, measure_id in measures:
    hosp_score = df_hcahps[
        (df_hcahps['Facility ID'] == hospital_id) &
        (df_hcahps['HCAHPS Measure ID'] == measure_id)
    ]['HCAHPS Answer Percent'].mean()
    state_avg = df_hcahps[
        (df_hcahps['State'] == hospital_state) &
        (df_hcahps['HCAHPS Measure ID'] == measure_id)
    ]['HCAHPS Answer Percent'].mean()
    nat_avg = df_hcahps[
        (df_hcahps['HCAHPS Measure ID'] == measure_id)
    ]['HCAHPS Answer Percent'].mean()
    comparison.append({
        'Measure': label,
        'Hospital': hosp_score,
        'State Avg': state_avg,
        'National Avg': nat_avg,
        'vs State': hosp_score - state_avg if pd.notnull(hosp_score) and pd.notnull(state_avg) else None,
        'vs National': hosp_score - nat_avg if pd.notnull(hosp_score) and pd.notnull(nat_avg) else None
    })

comp_df = pd.DataFrame(comparison)

# --- Display Table ---
st.subheader('Comparison Table')
st.dataframe(comp_df[['Measure', 'Hospital', 'State Avg', 'National Avg', 'vs State', 'vs National']].style.applymap(
    lambda v: 'background-color: #d4edda' if isinstance(v, (int, float)) and v > 0 else ('background-color: #f8d7da' if isinstance(v, (int, float)) and v < 0 else ''),
    subset=['vs State', 'vs National']
))

# --- Visualization ---
st.subheader('Benchmark Chart')
fig = go.Figure()

# Define consistent colors
colors = {
    'Hospital': 'rgb(0, 51, 102)',  # Dark blue
    'State Avg': 'rgb(128, 128, 128)',  # Gray
    'National Avg': 'rgb(173, 216, 230)'  # Light blue
}

# Add traces with consistent styling
fig.add_trace(go.Bar(
    x=comp_df['Measure'],
    y=comp_df['Hospital'],
    name='Hospital',
    marker_color=colors['Hospital'],
    text=comp_df['Hospital'].round(1).astype(str) + '%',
    textposition='outside'
))

fig.add_trace(go.Bar(
    x=comp_df['Measure'],
    y=comp_df['State Avg'],
    name='State Avg',
    marker_color=colors['State Avg'],
    text=comp_df['State Avg'].round(1).astype(str) + '%',
    textposition='outside'
))

fig.add_trace(go.Bar(
    x=comp_df['Measure'],
    y=comp_df['National Avg'],
    name='National Avg',
    marker_color=colors['National Avg'],
    text=comp_df['National Avg'].round(1).astype(str) + '%',
    textposition='outside'
))

# Update layout for better readability
fig.update_layout(
    barmode='group',
    xaxis_title='Measure',
    yaxis_title='Score (%)',
    plot_bgcolor='white',
    paper_bgcolor='white',
    showlegend=True,
    legend=dict(
        orientation='h',
        yanchor='bottom',
        y=1.02,
        xanchor='center',
        x=0.5
    ),
    margin=dict(t=100),  # Add top margin for value labels
    font=dict(size=14),
    xaxis=dict(
        tickfont=dict(size=12),
        title=dict(
            text='Measure',
            font=dict(size=14)
        )
    ),
    yaxis=dict(
        tickfont=dict(size=12),
        title=dict(
            text='Score (%)',
            font=dict(size=14)
        ),
        range=[0, 100]  # Set y-axis range to 0-100 for percentages
    )
)

# Remove chart border
fig.update_xaxes(showline=False, showgrid=False)
fig.update_yaxes(showline=False, showgrid=True, gridcolor='lightgray')

st.plotly_chart(fig, use_container_width=True)

# --- PowerPoint Export ---
def create_pptx(comp_df, hospital, chart_path):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # 6 is usually the fully blank layout
    slide = prs.slides.add_slide(blank_layout)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Remove any placeholder shapes (just in case)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    # Add title with padding
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), slide_width - Inches(1), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = f"HCAHPS Benchmark Report: {hospital}"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add date below title
    date_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), slide_width - Inches(1), Inches(0.4))
    date_frame = date_shape.text_frame
    date_frame.text = f"Analysis Date: {datetime.now().strftime('%Y-%m-%d')}"
    date_frame.paragraphs[0].font.size = Pt(18)
    date_frame.paragraphs[0].font.italic = True
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Table placement and sizing
    rows, cols = comp_df.shape
    table_width = slide_width - Inches(1)
    table_left = Inches(0.5)
    table_top = Inches(1.6)
    row_height = 0.18  # Shrink row height further
    table_height = Inches(0.4 + rows * row_height)
    table = slide.shapes.add_table(rows+1, cols, table_left, table_top, table_width, table_height).table

    # Set header style
    for j, col in enumerate(comp_df.columns):
        cell = table.cell(0, j)
        cell.text = col
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(221, 235, 247)  # light blue

    # Fill table data with right-aligned numbers
    for i in range(rows):
        for j in range(cols):
            val = comp_df.iloc[i, j]
            cell = table.cell(i+1, j)
            if pd.notnull(val):
                if isinstance(val, (int, float)):
                    cell.text = f"{val:.1f}" if isinstance(val, float) else str(val)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                else:
                    cell.text = str(val)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            else:
                cell.text = ""
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)

    # Auto-fit column widths
    for j in range(cols):
        max_len = max([len(str(comp_df.columns[j]))] + [len(str(comp_df.iloc[i, j])) for i in range(rows)])
        col_width = min(max(0.12 * max_len, 1.0), 2.2)
        table.columns[j].width = Inches(col_width)

    # Dynamically calculate chart position and size
    chart_margin = Inches(0.2)
    chart_top = table_top + table_height + chart_margin
    chart_height = slide_height - chart_top - Inches(0.3)  # leave some bottom margin
    chart_width = slide_width - Inches(1)
    chart_left = Inches(0.5)
    if chart_height < Inches(2):
        chart_height = Inches(2)  # minimum height

    # Create a new figure for the PowerPoint export with larger fonts and data labels
    export_fig = go.Figure()
    colors = {
        'Hospital': 'rgb(0, 51, 102)',
        'State Avg': 'rgb(128, 128, 128)',
        'National Avg': 'rgb(173, 216, 230)'
    }
    export_fig.add_trace(go.Bar(
        x=comp_df['Measure'],
        y=comp_df['Hospital'],
        name='Hospital',
        marker_color=colors['Hospital'],
        text=comp_df['Hospital'].round(1).astype(str) + '%',
        textposition='outside',
        textfont=dict(size=14)
    ))
    export_fig.add_trace(go.Bar(
        x=comp_df['Measure'],
        y=comp_df['State Avg'],
        name='State Avg',
        marker_color=colors['State Avg'],
        text=comp_df['State Avg'].round(1).astype(str) + '%',
        textposition='outside',
        textfont=dict(size=14)
    ))
    export_fig.add_trace(go.Bar(
        x=comp_df['Measure'],
        y=comp_df['National Avg'],
        name='National Avg',
        marker_color=colors['National Avg'],
        text=comp_df['National Avg'].round(1).astype(str) + '%',
        textposition='outside',
        textfont=dict(size=14)
    ))
    export_fig.update_layout(
        barmode='group',
        xaxis_title='Measure',
        yaxis_title='Score (%)',
        plot_bgcolor='white',
        paper_bgcolor='white',
        showlegend=True,
        legend=dict(
            orientation='h',
            yanchor='bottom',
            y=1.02,
            xanchor='center',
            x=0.5,
            font=dict(size=14)
        ),
        margin=dict(t=100),
        font=dict(size=16),
        xaxis=dict(
            tickfont=dict(size=16),  # Larger for PowerPoint
            title=dict(
                text='Measure',
                font=dict(size=18)
            )
        ),
        yaxis=dict(
            tickfont=dict(size=14),
            title=dict(
                text='Score (%)',
                font=dict(size=18)
            ),
            range=[0, 100]
        )
    )
    export_fig.update_xaxes(showline=False, showgrid=False)
    export_fig.update_yaxes(showline=False, showgrid=True, gridcolor='lightgray')
    export_fig.write_image(chart_path, width=1200, height=800)
    slide.shapes.add_picture(chart_path, chart_left, chart_top, chart_width, chart_height)
    return prs

def save_chart_as_image(fig):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
    fig.write_image(tmp.name)
    return tmp.name

if st.button('Download PowerPoint Report'):
    chart_img = save_chart_as_image(fig)
    pptx = create_pptx(comp_df, hospital, chart_img)
    tmp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    pptx.save(tmp_pptx.name)
    with open(tmp_pptx.name, 'rb') as f:
        st.download_button('Download Report', f, file_name=f"{hospital}_HCAHPS_Benchmark_{datetime.now().strftime('%Y%m%d')}.pptx")
    os.remove(chart_img)
    os.remove(tmp_pptx.name)
